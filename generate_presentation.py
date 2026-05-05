"""
SRUU Presentation - Modern Trendy Design (No Circles)
Style: Dark glassmorphism + sharp geometric cuts + neon accents + bold grid layout
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color Palette ─────────────────────────────────────────────────────────────
BG_DEEP      = RGBColor(0x07, 0x0E, 0x1A)   # near-black navy
BG_CARD      = RGBColor(0x0F, 0x1E, 0x33)   # dark card
BG_CARD2     = RGBColor(0x12, 0x25, 0x3D)   # slightly lighter card
PANEL_MID    = RGBColor(0x0A, 0x18, 0x2B)   # mid panel
NEON_CYAN    = RGBColor(0x00, 0xF5, 0xFF)   # neon cyan
NEON_ORANGE  = RGBColor(0xFF, 0x5C, 0x00)   # neon orange
NEON_GREEN   = RGBColor(0x00, 0xFF, 0xA3)   # neon green
NEON_PURPLE  = RGBColor(0xAA, 0x00, 0xFF)   # neon purple
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT    = RGBColor(0xB0, 0xC8, 0xE0)
DIM_TEXT     = RGBColor(0x60, 0x88, 0xAA)
STRIPE       = RGBColor(0x1A, 0x35, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Core drawing helpers ──────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line=None, lw=None):
    sh = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw or 1)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return tb

def mtxt(slide, lines, l, t, w, h, size=12, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb

def bg(slide):
    """Full dark background"""
    rect(slide, 0, 0, 13.33, 7.5, fill=BG_DEEP)

def neon_line(slide, l, t, w, color, h=0.04):
    """Thin glowing accent line - h must be a number"""
    rect(slide, l, t, w, float(h), fill=color)

def tag(slide, text, l, t, color, text_color=None):
    """Pill-shaped label"""
    rect(slide, l, t, len(text)*0.11 + 0.3, 0.32, fill=color)
    txt(slide, text, l+0.1, t+0.02, len(text)*0.11+0.1, 0.28,
        size=9, bold=True, color=text_color or BG_DEEP, align=PP_ALIGN.LEFT)

def card(slide, l, t, w, h, accent=None, accent_w=0.055):
    """Dark card with optional left accent bar"""
    rect(slide, l, t, w, h, fill=BG_CARD)
    if accent:
        rect(slide, l, t, accent_w, h, fill=accent)

def glass_card(slide, l, t, w, h, border_color=None):
    """Glassmorphism-style card"""
    rect(slide, l, t, w, h, fill=BG_CARD2)
    if border_color:
        rect(slide, l, t, w, 0.035, fill=border_color)
        rect(slide, l, t, 0.035, h, fill=border_color)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE  (Split diagonal layout)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

# Left dark panel
rect(slide, 0, 0, 7.8, 7.5, fill=PANEL_MID)

# Right neon panel — sharp diagonal via overlapping rects (staircase illusion)
rect(slide, 7.8, 0, 5.53, 7.5, fill=BG_DEEP)
# Diagonal accent strip — simulated with thin slanted rects
for i in range(18):
    rect(slide, 7.55 + i*0.015, i*0.42, 0.22, 7.5, fill=BG_CARD2)

# Top neon bar
neon_line(slide, 0, 0, 13.33, NEON_CYAN, 0.055)

# Bottom neon strip
neon_line(slide, 0, 7.4, 13.33, NEON_ORANGE, 0.065)
txt(slide, "Sara Arkoub  |  UMBB M1 Genie Logiciel  |  2025-2026",
    0.5, 7.42, 9, 0.28, size=9, color=DIM_TEXT)

# ── Left content ──
txt(slide, "PROJET SMA", 0.6, 0.25, 5, 0.35, size=10, bold=True, color=NEON_CYAN)
neon_line(slide, 0.6, 0.65, 1.5, NEON_ORANGE, 0.04)

txt(slide, "SRUU", 0.6, 0.78, 7, 1.6, size=88, bold=True, color=WHITE)

txt(slide, "Systeme de Reponse", 0.6, 2.42, 7, 0.5, size=22, bold=False, color=GRAY_TEXT)
txt(slide, "aux Urgences Urbaines", 0.6, 2.88, 7, 0.5, size=22, bold=True, color=NEON_CYAN)

neon_line(slide, 0.6, 3.45, 5.5, NEON_CYAN, 0.03)

txt(slide, "Urban Emergency Response System", 0.6, 3.58, 7, 0.4,
    size=14, italic=True, color=GRAY_TEXT)
txt(slide, "Multi-Agent System with JADE Framework", 0.6, 3.98, 7, 0.35,
    size=12, color=DIM_TEXT)

# Stat pills
stats = [("9", "AGENTS"), ("JADE", "FRAMEWORK"), ("FIPA-ACL", "PROTOCOL")]
for i, (val, lab) in enumerate(stats):
    lx = 0.6 + i * 2.4
    rect(slide, lx, 4.65, 2.1, 0.75, fill=STRIPE)
    neon_line(slide, lx, 4.65, 2.1, NEON_CYAN if i == 0 else NEON_ORANGE if i==1 else NEON_GREEN, 0.04)
    txt(slide, val, lx + 0.15, 4.72, 1.8, 0.35, size=18, bold=True, color=WHITE)
    txt(slide, lab, lx + 0.15, 5.07, 1.8, 0.25, size=8, bold=True, color=NEON_CYAN if i==0 else NEON_ORANGE if i==1 else NEON_GREEN)

# Emergency type tags
for i, (em, col) in enumerate([("FIRE", NEON_ORANGE), ("MEDICAL", NEON_GREEN), ("BIOHAZARD", NEON_PURPLE)]):
    lx = 0.6 + i * 2.3
    rect(slide, lx, 5.7, 2.1, 0.45, fill=BG_CARD, line=col, lw=1.2)
    txt(slide, em, lx+0.15, 5.77, 1.8, 0.32, size=11, bold=True, color=col)

# ── Right panel ──
txt(slide, "SYSTEME", 8.2, 1.2, 5, 0.8, size=38, bold=True, color=DIM_TEXT)
txt(slide, "MULTI-AGENT", 8.2, 1.9, 5, 0.8, size=28, bold=True, color=NEON_CYAN)

# Vertical feature list
features = [
    (NEON_ORANGE, "Coordination temps reel"),
    (NEON_GREEN,  "Protocole Contract Net"),
    (NEON_CYAN,   "Ontologie FIPA-ACL"),
    (NEON_PURPLE, "9 agents autonomes"),
]
for i, (col, feat) in enumerate(features):
    ty = 3.0 + i * 0.82
    rect(slide, 8.2, ty, 0.06, 0.55, fill=col)
    txt(slide, feat, 8.45, ty + 0.06, 4.6, 0.42, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE  (Numbered grid layout)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
neon_line(slide, 0, 0, 13.33, NEON_CYAN, 0.05)
neon_line(slide, 0, 7.44, 13.33, NEON_PURPLE, 0.06)

txt(slide, "SOMMAIRE", 0.6, 0.15, 4, 0.35, size=10, bold=True, color=NEON_CYAN)
txt(slide, "Plan de la presentation", 0.6, 0.5, 10, 0.65, size=34, bold=True, color=WHITE)
neon_line(slide, 0.6, 1.2, 3.5, NEON_ORANGE, 0.04)

items = [
    ("01", "Introduction",               "Contexte et objectifs",                  NEON_CYAN),
    ("02", "Problematique",              "Defis de la gestion des urgences",       NEON_ORANGE),
    ("03", "Solution Proposee",          "Architecture multi-agent SRUU",          NEON_GREEN),
    ("04", "Les Agents",                 "9 types d'agents autonomes",             NEON_PURPLE),
    ("05", "Communication",              "FIPA-ACL, Contract Net, Ontologie",      NEON_CYAN),
    ("06", "Scenarios & Resultats",      "Simulation et metriques",                NEON_ORANGE),
    ("07", "Conclusion",                 "Bilan et perspectives",                  NEON_GREEN),
]

# 3+4 grid layout
positions = [
    (0.5,  1.45), (4.65, 1.45), (8.8, 1.45),
    (0.5,  3.85), (4.65, 3.85), (8.8, 3.85), (4.65, 6.1),
]
for i, ((lx, ty), (num, title, sub, col)) in enumerate(zip(positions, items)):
    w = 3.8
    rect(slide, lx, ty, w, 2.0, fill=BG_CARD)
    neon_line(slide, lx, ty, w, col, 0.05)
    # Big number
    txt(slide, num, lx+0.18, ty+0.1, 1.2, 0.8, size=36, bold=True, color=col)
    neon_line(slide, lx+0.18, ty+0.95, 3.4, STRIPE, 0.025)
    txt(slide, title, lx+0.18, ty+1.07, 3.5, 0.42, size=14, bold=True, color=WHITE)
    txt(slide, sub,   lx+0.18, ty+1.52, 3.5, 0.35, size=9,  color=DIM_TEXT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION  (Split panel + icon grid)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
neon_line(slide, 0, 0, 13.33, NEON_GREEN, 0.05)
neon_line(slide, 0, 7.44, 13.33, NEON_GREEN, 0.05)

# Left dark column
rect(slide, 0, 0.05, 4.2, 7.4, fill=PANEL_MID)
rect(slide, 4.2, 0.05, 0.04, 7.4, fill=NEON_GREEN)

txt(slide, "01", 0.35, 0.2, 1.5, 0.8, size=52, bold=True, color=STRIPE)
txt(slide, "INTRODUCTION", 0.35, 0.95, 3.7, 0.38, size=10, bold=True, color=NEON_GREEN)
neon_line(slide, 0.35, 1.38, 3.5, NEON_GREEN, 0.04)
txt(slide, "Contexte\ndu Projet", 0.35, 1.55, 3.7, 1.0, size=26, bold=True, color=WHITE)

# Stat blocks
stats_l = [("9", "Agents autonomes"), ("3", "Types d'incidents"), ("JADE", "Java Framework"), ("100x100", "Grille simulation")]
for i, (val, lab) in enumerate(stats_l):
    ty = 2.85 + i * 1.1
    rect(slide, 0.35, ty, 3.55, 0.88, fill=BG_CARD)
    rect(slide, 0.35, ty, 0.05, 0.88, fill=NEON_GREEN)
    txt(slide, val, 0.6, ty+0.06, 1.5, 0.45, size=24, bold=True, color=NEON_GREEN)
    txt(slide, lab, 0.6, ty+0.52, 3.1, 0.3, size=10, color=GRAY_TEXT)

# Right content
txt(slide, "Qu'est-ce que le SRUU ?", 4.5, 0.18, 8.5, 0.55, size=22, bold=True, color=WHITE)
neon_line(slide, 4.5, 0.78, 8.5, NEON_GREEN, 0.04)

intro = (
    "Le Systeme de Reponse aux Urgences Urbaines (SRUU) est une simulation intelligente "
    "de coordination d'urgences urbaines basee sur les Systemes Multi-Agents (SMA). "
    "Developpe avec JADE (Java Agent Development Framework), il orchestre 9 agents "
    "autonomes qui collaborent via le protocole Contract Net pour detecter, evaluer "
    "et resoudre des incidents en temps reel."
)
txt(slide, intro, 4.5, 0.95, 8.5, 1.8, size=12, color=GRAY_TEXT)

# Context cards — sharp left-bordered
ctx = [
    (NEON_CYAN,   "Contexte Urbain",
     "Les villes modernes font face a des crises complexes necessitant une coordination rapide entre services heterogenes."),
    (NEON_ORANGE, "Approche SMA",
     "Les SMA offrent une solution decentralisee, reactive et robuste via des agents autonomes specialises."),
    (NEON_GREEN,  "Objectif",
     "Simuler une reponse coordonnee et optimale via Contract Net avec une fonction d'utilite multi-criteres."),
    (NEON_PURPLE, "Technologies",
     "JADE 4.5 + Java 11 + Maven + FIPA-ACL + Ontologie XML + Rapports JSON automatiques."),
]
for i, (col, title, body) in enumerate(ctx):
    ty = 2.95 + i * 1.1
    rect(slide, 4.5, ty, 8.5, 0.92, fill=BG_CARD)
    rect(slide, 4.5, ty, 0.06, 0.92, fill=col)
    txt(slide, title, 4.75, ty+0.06, 8.1, 0.35, size=12, bold=True, color=col)
    txt(slide, body,  4.75, ty+0.45, 8.1, 0.45, size=10, color=GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLÉMATIQUE  (Hexagon-inspired card grid, no circles)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
neon_line(slide, 0, 0, 13.33, NEON_ORANGE, 0.05)
neon_line(slide, 0, 7.44, 13.33, NEON_ORANGE, 0.05)

txt(slide, "02  PROBLEMATIQUE", 0.6, 0.12, 8, 0.38, size=10, bold=True, color=NEON_ORANGE)
txt(slide, "Defis de la Gestion des Urgences Urbaines", 0.6, 0.48, 12.5, 0.62,
    size=28, bold=True, color=WHITE)
neon_line(slide, 0.6, 1.12, 12.1, STRIPE, 0.03)

# Central question block
rect(slide, 0.5, 1.22, 12.3, 0.92, fill=BG_CARD)
neon_line(slide, 0.5, 1.22, 12.3, NEON_ORANGE, 0.05)
txt(slide, "Comment coordonner efficacement des ressources heterogenes (ambulances, pompiers, police) "
    "en temps reel, face a des incidents multiples et simultanees dans un environnement urbain ?",
    0.75, 1.3, 11.8, 0.78, size=12, italic=True, color=WHITE)

# Problem cards 3x2 — diamond-corner style (sharp rectangles with diagonal top bar)
problems = [
    (NEON_ORANGE, "Temps Reel",
     "Decisions en < 500 ms pour eviter la degradation — chaque seconde compte dans une urgence urbaine."),
    (NEON_CYAN,   "Coordination",
     "Unites heterogenes collaborent sans conflit — partage de ressources et evitement de doublons."),
    (NEON_GREEN,  "Localisation Optimale",
     "Assigner l'unite la plus proche ET la plus adaptee parmi N candidats distribues."),
    (NEON_ORANGE, "Pannes & Ressources",
     "Gestion des defaillances : camion a court d'eau → reassignation automatique immediate."),
    (NEON_CYAN,   "Communication Fiable",
     "Echange de messages structures (FIPA-ACL) sans deadlock ni perte de donnees."),
    (NEON_GREEN,  "Decision Optimale",
     "Selectionner la meilleure unite via une fonction d'utilite multi-criteres mathematique."),
]
for i, (col, title, body) in enumerate(problems):
    c = i % 3; r = i // 3
    lx = 0.5 + c * 4.3; ty = 2.42 + r * 2.4
    w, h = 4.0, 2.1
    rect(slide, lx, ty, w, h, fill=BG_CARD)
    # Top accent bar
    neon_line(slide, lx, ty, w, col, 0.055)
    # Number watermark
    txt(slide, f"0{i+1}", lx+2.8, ty+0.1, 1.1, 0.8, size=42, bold=True, color=STRIPE)
    txt(slide, title, lx+0.2, ty+0.18, 3.6, 0.42, size=13, bold=True, color=col)
    neon_line(slide, lx+0.2, ty+0.65, 3.6, STRIPE, 0.025)
    txt(slide, body, lx+0.2, ty+0.78, 3.6, 1.2, size=10, color=GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SOLUTION / ARCHITECTURE  (Layered pipeline)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
neon_line(slide, 0, 0, 13.33, NEON_CYAN, 0.05)
neon_line(slide, 0, 7.44, 13.33, NEON_CYAN, 0.05)

txt(slide, "03  SOLUTION PROPOSEE", 0.6, 0.12, 8, 0.35, size=10, bold=True, color=NEON_CYAN)
txt(slide, "Architecture du Systeme SRUU", 0.6, 0.45, 10, 0.62, size=28, bold=True, color=WHITE)
neon_line(slide, 0.6, 1.1, 6, NEON_CYAN, 0.04)

layers = [
    (NEON_GREEN,  "PERCEPTION",  "Sensor Agents x3",
     "Detectent les incidents sur grille 100x100 avec probabilite 30% par cycle de 5s."),
    (NEON_CYAN,   "DECISION",    "Dispatcher Agent x1",
     "Evalue les propositions via fonction d'utilite U(unit) et assigne via Contract Net."),
    (NEON_ORANGE, "ACTION",      "Ambulance x2  |  FireTruck x2  |  Police x2  |  Biohazard x1",
     "Repondent aux incidents, gerent leurs ressources, signalent les pannes (ABORT)."),
    (NEON_PURPLE, "SUPPORT",     "MedicalCoordinator  |  TrafficController  |  Logger",
     "Gestion hospitaliere, optimisation routes d'urgence, audit JSON automatique."),
]

arrow_col = STRIPE
for i, (col, layer, title, desc) in enumerate(layers):
    ty = 1.3 + i * 1.45
    # Arrow connector between layers
    if i > 0:
        rect(slide, 0.85, ty - 0.28, 0.05, 0.28, fill=col)
        # Arrowhead simulation
        rect(slide, 0.78, ty - 0.08, 0.18, 0.04, fill=col)
    # Layer block
    rect(slide, 0.6, ty, 12.3, 1.2, fill=BG_CARD)
    neon_line(slide, 0.6, ty, 12.3, col, 0.055)
    # Layer badge
    rect(slide, 0.6, ty+0.12, 1.9, 0.36, fill=col)
    txt(slide, layer, 0.65, ty+0.14, 1.85, 0.32, size=9, bold=True,
        color=BG_DEEP, align=PP_ALIGN.CENTER)
    txt(slide, title, 2.7, ty+0.1, 9.8, 0.4, size=14, bold=True, color=WHITE)
    neon_line(slide, 2.7, ty+0.55, 9.8, STRIPE, 0.025)
    txt(slide, desc, 2.7, ty+0.65, 9.8, 0.5, size=10.5, color=GRAY_TEXT)

# Tech stack strip
neon_line(slide, 0, 7.1, 13.33, STRIPE, 0.04)
rect(slide, 0, 7.14, 13.33, 0.3, fill=BG_CARD)
txt(slide, "Java 11+  |  JADE 4.5.0  |  Maven 3.6+  |  FIPA-ACL  |  XML Ontology  |  JSON Reports",
    0.5, 7.16, 12.5, 0.28, size=10, bold=True, color=NEON_CYAN, align=PP_ALIGN.CENTER)

# Architecture diagram right annotation
rect(slide, 9.8, 1.25, 3.1, 5.7, fill=PANEL_MID)
neon_line(slide, 9.8, 1.25, 3.1, STRIPE, 0.035)
txt(slide, "STACK TECHNIQUE", 9.95, 1.38, 2.8, 0.35, size=9, bold=True, color=DIM_TEXT)
tech = [
    (NEON_CYAN,   "JADE 4.5.0"),
    (NEON_GREEN,  "Java 11+"),
    (NEON_ORANGE, "Maven Build"),
    (NEON_PURPLE, "FIPA-ACL"),
    (NEON_CYAN,   "Contract Net"),
    (NEON_GREEN,  "FSM Behaviours"),
    (NEON_ORANGE, "XML Ontology"),
    (NEON_PURPLE, "JSON Reports"),
]
for i, (col, t) in enumerate(tech):
    ty = 1.85 + i * 0.58
    rect(slide, 9.95, ty, 0.045, 0.38, fill=col)
    txt(slide, t, 10.1, ty+0.04, 2.65, 0.32, size=11, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LES AGENTS  (3x3 grid sharp cards)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
neon_line(slide, 0, 0, 13.33, NEON_GREEN, 0.05)
neon_line(slide, 0, 7.44, 13.33, NEON_GREEN, 0.05)

txt(slide, "04  LES AGENTS DU SYSTEME", 0.5, 0.12, 9, 0.35, size=10, bold=True, color=NEON_GREEN)
txt(slide, "9 Types d'Agents Autonomes", 0.5, 0.45, 10, 0.62, size=28, bold=True, color=WHITE)
neon_line(slide, 0.5, 1.1, 5, NEON_GREEN, 0.04)

agents = [
    ("DISPATCHER",    "x1", NEON_CYAN,   "Coordinateur central. Protocole Contract Net. Selection via U(unit)."),
    ("SENSOR",        "x3", NEON_GREEN,  "Detection incidents. Grille 100x100. Probabilite 30% / 5s cycle."),
    ("AMBULANCE",     "x2", NEON_GREEN,  "MEDICAL. Vitesse 2u/tick. Coordination hospitaliere auto."),
    ("FIRE TRUCK",    "x2", NEON_ORANGE, "FIRE + BIOHAZARD. Vitesse 3u/tick. Reservoir 100L eau."),
    ("POLICE UNIT",   "x2", NEON_CYAN,   "PERIMETER. Vitesse 3u/tick. Controle foule & perimetre."),
    ("BIOHAZARD CU",  "x1", NEON_PURPLE, "HAZMAT. Vitesse 2u/tick. Confinement matières dangereuses."),
    ("MED COORD.",    "x1", NEON_GREEN,  "Gestion capacite hopitaux. Allocation lits. Transferts patients."),
    ("TRAFFIC CTRL",  "x1", NEON_ORANGE, "Optimisation routes urgence. Corridors prioritaires GPS."),
    ("LOGGER",        "x1", NEON_CYAN,   "Audit complet. Rapports JSON. KPI: temps reponse, incidents."),
]

cols_n = 3
for i, (name, count, col, desc) in enumerate(agents):
    c = i % cols_n; r = i // cols_n
    lx = 0.4 + c * 4.3
    ty = 1.32 + r * 2.0
    w, h = 4.0, 1.75
    rect(slide, lx, ty, w, h, fill=BG_CARD)
    neon_line(slide, lx, ty, w, col, 0.055)
    # Name + count row
    txt(slide, name, lx+0.18, ty+0.12, 3.0, 0.4, size=13, bold=True, color=WHITE)
    rect(slide, lx+3.1, ty+0.1, 0.72, 0.35, fill=col)
    txt(slide, count, lx+3.1, ty+0.1, 0.72, 0.35, size=11, bold=True,
        color=BG_DEEP, align=PP_ALIGN.CENTER)
    neon_line(slide, lx+0.18, ty+0.58, 3.6, STRIPE, 0.025)
    txt(slide, desc, lx+0.18, ty+0.72, 3.65, 0.95, size=10, color=GRAY_TEXT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — COMMUNICATION & PROTOCOLES  (Flow diagram + code block)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
neon_line(slide, 0, 0, 13.33, NEON_PURPLE, 0.05)
neon_line(slide, 0, 7.44, 13.33, NEON_PURPLE, 0.05)

txt(slide, "05  COMMUNICATION & PROTOCOLES", 0.6, 0.12, 9, 0.35, size=10, bold=True, color=NEON_PURPLE)
txt(slide, "FIPA-ACL  |  Contract Net Protocol", 0.6, 0.45, 11, 0.62, size=28, bold=True, color=WHITE)
neon_line(slide, 0.6, 1.1, 7, NEON_PURPLE, 0.04)

# Contract Net flow — horizontal pipeline
flow = [
    ("INCIDENT\nDETECTE",  NEON_ORANGE, "Sensor→Dispatcher\nINFORM"),
    ("CFP\nENVOYE",        NEON_CYAN,   "Dispatcher→Units\nCALL FOR PROP."),
    ("PROPOSE\nRECU",      NEON_GREEN,  "Units→Dispatcher\nSCORE + DETAILS"),
    ("ACCEPT /\nREJECT",   NEON_ORANGE, "Best Unit\nACCEPT-PROPOSAL"),
    ("MISSION\nTERMINEE",  NEON_GREEN,  "Unit→Dispatcher\nINFORM DONE"),
]
bw = 2.0
total = len(flow)*bw + (len(flow)-1)*0.35
sx = (13.33 - total) / 2

for i, (step, col, detail) in enumerate(flow):
    lx = sx + i*(bw + 0.35)
    # Arrow
    if i > 0:
        rect(slide, lx-0.35, 2.12, 0.35, 0.04, fill=col)
        rect(slide, lx-0.09, 2.05, 0.04, 0.18, fill=col)
    rect(slide, lx, 1.32, bw, 1.62, fill=BG_CARD)
    neon_line(slide, lx, 1.32, bw, col, 0.055)
    # Step number badge
    rect(slide, lx + bw - 0.42, 1.38, 0.35, 0.35, fill=col)
    txt(slide, str(i+1), lx+bw-0.42, 1.38, 0.35, 0.35, size=12, bold=True,
        color=BG_DEEP, align=PP_ALIGN.CENTER)
    txt(slide, step, lx+0.12, 1.78, bw-0.15, 0.72, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, detail, lx+0.08, 2.5, bw-0.1, 0.42, size=8.5, color=DIM_TEXT, align=PP_ALIGN.CENTER)

# Utility function block — full width
rect(slide, 0.5, 3.12, 12.3, 1.05, fill=BG_CARD)
neon_line(slide, 0.5, 3.12, 12.3, NEON_ORANGE, 0.055)
rect(slide, 0.5, 3.12, 2.4, 0.38, fill=NEON_ORANGE)
txt(slide, "FONCTION D'UTILITE", 0.6, 3.14, 2.3, 0.34, size=9, bold=True,
    color=BG_DEEP, align=PP_ALIGN.CENTER)
txt(slide,
    "U(unit)  =  0.35 x ProxScore  +  0.25 x TypeMatch  +  0.20 x StatusScore  +  0.20 x SeverityScore",
    0.7, 3.6, 12.0, 0.45, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "L'unite avec le score MAX est selectionnee  |  Distance euclidienne sur grille 100x100",
    0.7, 4.05, 12.0, 0.3, size=10, color=DIM_TEXT, align=PP_ALIGN.CENTER)

# ACL message block
rect(slide, 0.5, 4.42, 5.8, 2.72, fill=RGBColor(0x04, 0x0C, 0x18))
neon_line(slide, 0.5, 4.42, 5.8, NEON_GREEN, 0.04)
rect(slide, 0.5, 4.42, 2.2, 0.32, fill=NEON_GREEN)
txt(slide, "MESSAGE FIPA-ACL", 0.56, 4.44, 2.15, 0.28, size=8, bold=True,
    color=BG_DEEP, align=PP_ALIGN.CENTER)
acl = [
    ":performative   PROPOSE",
    ":sender         ambulance-agent-1",
    ":receiver       dispatcher-agent",
    ":content        { score:0.82, dist:12.5, type:MEDICAL }",
    ":ontology       emergency-response-ontology",
    ":language       FIPA-SL",
    ":conversation-id  INC-20260429-001",
]
mtxt(slide, acl, 0.65, 4.82, 5.5, 2.2, size=10, color=NEON_GREEN)

# Protocols list
rect(slide, 6.7, 4.42, 6.1, 2.72, fill=BG_CARD)
neon_line(slide, 6.7, 4.42, 6.1, NEON_CYAN, 0.04)
txt(slide, "STANDARDS & PROTOCOLES", 6.85, 4.5, 5.8, 0.35, size=10, bold=True, color=NEON_CYAN)
neon_line(slide, 6.85, 4.88, 5.8, STRIPE, 0.025)
protos = [
    (NEON_PURPLE, "Contract Net Protocol (FIPA)",   "Negociation distribuee multi-agents"),
    (NEON_CYAN,   "FIPA-ACL Messaging",             "Communication formelle structuree"),
    (NEON_ORANGE, "XML Ontologie",                  "Semantique partagee inter-agents"),
    (NEON_GREEN,  "FSM Behaviours JADE",            "Automate d'etat fini par agent"),
    (NEON_PURPLE, "Hierarchical Dispatcher",        "Architecture centralisee optimisee"),
]
for i, (col, p, d) in enumerate(protos):
    ty = 5.0 + i*0.42
    rect(slide, 6.85, ty, 0.04, 0.28, fill=col)
    txt(slide, p, 7.05, ty+0.01, 3.0, 0.28, size=10, bold=True, color=WHITE)
    txt(slide, d, 10.1, ty+0.01, 2.55, 0.28, size=9, color=DIM_TEXT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SCÉNARIOS & RÉSULTATS  (Timeline + metrics dashboard)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
neon_line(slide, 0, 0, 13.33, NEON_ORANGE, 0.05)
neon_line(slide, 0, 7.44, 13.33, NEON_ORANGE, 0.05)

txt(slide, "06  SCENARIOS & RESULTATS", 0.5, 0.12, 9, 0.35, size=10, bold=True, color=NEON_ORANGE)
txt(slide, "Simulation & Metriques de Performance", 0.5, 0.45, 11, 0.62, size=28, bold=True, color=WHITE)
neon_line(slide, 0.5, 1.1, 7, NEON_ORANGE, 0.04)

# 3 scenario cards
scenarios = [
    ("01", "Incidents Multi-Types", NEON_GREEN,
     "Declenchement FIRE + MEDICAL + BIOHAZARD simultanement.\n"
     "Selection correcte du type d'unite via U(unit).\n"
     "Verification temps reponse & matching de type."),
    ("02", "Incidents Simultanees", NEON_CYAN,
     "9 incidents generes en parallele (3 par Sensor).\n"
     "Traitement parallele sans deadlock ni blocage.\n"
     "Deploiement simultane de plusieurs unites."),
    ("03", "Panne de Ressource", NEON_ORANGE,
     "Depletion eau camion → message ABORT automatique.\n"
     "Reassignation vers unite disponible en < 1s.\n"
     "Test robustesse et recuperation systeme."),
]
for i, (num, title, col, body) in enumerate(scenarios):
    ty = 1.32
    lx = 0.45 + i * 4.28
    w, h = 3.95, 3.3
    rect(slide, lx, ty, w, h, fill=BG_CARD)
    neon_line(slide, lx, ty, w, col, 0.055)
    # Big number watermark
    txt(slide, num, lx+2.5, ty+0.08, 1.3, 0.9, size=52, bold=True, color=STRIPE)
    txt(slide, title, lx+0.18, ty+0.18, 3.6, 0.5, size=14, bold=True, color=col)
    neon_line(slide, lx+0.18, ty+0.74, 3.55, STRIPE, 0.025)
    mtxt(slide, body.split("\n"), lx+0.18, ty+0.9, 3.6, 2.2, size=10.5, color=GRAY_TEXT)

# JSON sample strip
rect(slide, 0.45, 4.78, 12.4, 0.5, fill=RGBColor(0x04, 0x0C, 0x18))
neon_line(slide, 0.45, 4.78, 12.4, NEON_GREEN, 0.035)
txt(slide,
    '{ "status":"COMPLETED", "assignedUnit":"FireTruck-1", "responseTime":45.2, "severity":8, "type":"FIRE" }',
    0.65, 4.85, 12.0, 0.38, size=10, color=NEON_GREEN)

# Metrics dashboard — bottom
metrics = [
    ("< 5s",     "Demarrage",         NEON_GREEN),
    ("100-500ms","Decision Dispatch", NEON_CYAN),
    ("45-60s",   "Temps Reponse",     NEON_ORANGE),
    ("3-5/min",  "Incidents/min",     NEON_GREEN),
    ("~100 MB",  "Memoire",           NEON_CYAN),
    ("0",        "Deadlocks",         NEON_GREEN),
]
mw = 12.4 / len(metrics)
for i, (val, label, col) in enumerate(metrics):
    lx = 0.45 + i * mw
    rect(slide, lx, 5.45, mw-0.08, 1.72, fill=BG_CARD)
    neon_line(slide, lx, 5.45, mw-0.08, col, 0.05)
    txt(slide, val, lx+0.15, 5.58, mw-0.3, 0.62, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    neon_line(slide, lx+0.12, 6.2, mw-0.28, STRIPE, 0.025)
    txt(slide, label, lx+0.1, 6.28, mw-0.2, 0.62, size=9, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSION  (Bold summary + perspective grid)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
neon_line(slide, 0, 0, 13.33, NEON_GREEN, 0.05)
neon_line(slide, 0, 7.44, 13.33, NEON_CYAN, 0.05)

# Right accent panel
rect(slide, 9.2, 0.05, 4.13, 7.4, fill=PANEL_MID)
rect(slide, 9.2, 0.05, 0.04, 7.4, fill=NEON_GREEN)

txt(slide, "07  CONCLUSION", 0.5, 0.12, 6, 0.35, size=10, bold=True, color=NEON_GREEN)
txt(slide, "Bilan &\nPerspectives", 0.5, 0.45, 8.5, 1.3, size=38, bold=True, color=WHITE)
neon_line(slide, 0.5, 1.82, 5, NEON_GREEN, 0.04)

# Achievements
ach_title = "Ce qui a ete realise"
txt(slide, ach_title, 0.5, 2.0, 5, 0.4, size=13, bold=True, color=NEON_GREEN)
achievements = [
    "9 agents autonomes avec comportements FSM complets",
    "Protocole Contract Net FIPA de bout en bout",
    "Fonction d'utilite multi-criteres (4 poids)",
    "Gestion robuste des pannes + reassignation auto",
    "Rapports JSON complets avec metriques KPI",
    "Zero deadlock sur tous les scenarios de test",
]
for i, ach in enumerate(achievements):
    ty = 2.5 + i * 0.7
    rect(slide, 0.5, ty, 8.5, 0.58, fill=BG_CARD)
    rect(slide, 0.5, ty, 0.05, 0.58, fill=NEON_GREEN)
    txt(slide, ach, 0.72, ty+0.1, 8.1, 0.38, size=11, color=GRAY_TEXT)

# Right panel — key numbers + summary
txt(slide, "EN CHIFFRES", 9.45, 0.25, 3.7, 0.35, size=10, bold=True, color=DIM_TEXT)
kn = [
    ("9",    "Agents",       NEON_GREEN),
    ("3",    "Scenarios",    NEON_CYAN),
    ("0",    "Deadlocks",    NEON_GREEN),
    ("100%", "FIPA-ACL",     NEON_ORANGE),
]
for i, (val, lab, col) in enumerate(kn):
    c = i % 2; r = i // 2
    lx = 9.45 + c * 1.9; ty = 0.7 + r * 1.35
    rect(slide, lx, ty, 1.68, 1.12, fill=BG_CARD)
    neon_line(slide, lx, ty, 1.68, col, 0.05)
    txt(slide, val, lx+0.12, ty+0.1, 1.45, 0.55, size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(slide, lab, lx+0.12, ty+0.68, 1.45, 0.3, size=9, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

neon_line(slide, 9.35, 3.18, 3.7, STRIPE, 0.03)
txt(slide, "PERSPECTIVES", 9.45, 3.3, 3.7, 0.35, size=10, bold=True, color=DIM_TEXT)

persp = [
    (NEON_CYAN,   "Dashboard Web",       "Interface temps reel & carte interactive"),
    (NEON_ORANGE, "ML Optimization",     "Apprentissage poids fonction d'utilite"),
    (NEON_GREEN,  "Multi-Dispatcher",    "Coordination hierarchique inter-villes"),
    (NEON_PURPLE, "Predictive Model",    "Anticipation incidents via historique"),
]
for i, (col, p, d) in enumerate(persp):
    ty = 3.75 + i * 0.88
    rect(slide, 9.35, ty, 3.7, 0.72, fill=BG_CARD)
    rect(slide, 9.35, ty, 0.05, 0.72, fill=col)
    txt(slide, p, 9.55, ty+0.06, 3.4, 0.3, size=11, bold=True, color=col)
    txt(slide, d, 9.55, ty+0.38, 3.4, 0.28, size=9, color=GRAY_TEXT)

# Bottom thank you
rect(slide, 0, 7.1, 13.33, 0.35, fill=BG_CARD)
neon_line(slide, 0, 7.1, 13.33, NEON_GREEN, 0.035)
txt(slide, "Merci pour votre attention", 0.5, 7.14, 7, 0.28, size=13, bold=True, color=WHITE)
txt(slide, "Sara Arkoub  |  UMBB M1 Genie Logiciel  |  2025-2026",
    7.5, 7.14, 5.8, 0.28, size=11, color=NEON_CYAN, align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
out = r"c:\Users\Administrator\Desktop\projet_sma\SRUU_Presentation_v2.pptx"
prs.save(out)
print("Generated:", out, "| Slides:", len(prs.slides))
